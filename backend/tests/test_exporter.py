import zipfile

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.util import Inches

from app.exporter import _bring_text_shapes_to_front, export_pptx
from app.models import BBox, Component, Project
from app.settings import Settings
from app.storage import ProjectStore


def project_store(tmp_path):
    settings = Settings(data_dir=tmp_path / "data")
    return ProjectStore(settings)


def source_slide(tmp_path):
    image_path = tmp_path / "slide.png"
    image = Image.new("RGB", (160, 90), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([12, 12, 148, 30], fill="#dbeafe")
    draw.rectangle([20, 17, 70, 24], fill="black")
    draw.ellipse([105, 45, 135, 75], fill="#d71920")
    image.save(image_path)
    return image_path


def shape_counts(pptx_path):
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.has_text_frame and shape.text]
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]
    return len(text_shapes), len(picture_shapes)


def media_images(pptx_path):
    images = []
    with zipfile.ZipFile(pptx_path) as package:
        for name in sorted(item for item in package.namelist() if item.startswith("ppt/media/")):
            with package.open(name) as handle:
                images.append(Image.open(handle).convert("RGBA").copy())
    return images


def non_empty_text_order(pptx_path):
    prs = Presentation(str(pptx_path))
    order = []
    for index, shape in enumerate(prs.slides[0].shapes):
        text = ""
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
        order.append((index, bool(text), str(shape.shape_type), text))
    return order


def assert_text_shapes_are_frontmost(pptx_path):
    order = non_empty_text_order(pptx_path)
    first_text = min((index for index, has_text, _, _ in order if has_text), default=None)
    last_visual = max((index for index, has_text, _, _ in order if not has_text), default=None)
    assert first_text is not None
    assert last_visual is not None
    assert first_text > last_visual


def test_bring_text_shapes_to_front_moves_existing_text_after_later_shapes():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(2.0), Inches(0.4))
    textbox.text_frame.paragraphs[0].text = "Front text"
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.1), Inches(0.1), Inches(3.0), Inches(0.8))

    assert not slide.shapes[-1].has_text_frame or slide.shapes[-1].text != "Front text"

    _bring_text_shapes_to_front(slide)

    assert slide.shapes[-1].has_text_frame
    assert slide.shapes[-1].text == "Front text"


def test_export_reconstructs_slide_even_when_analysis_only_has_text(tmp_path):
    image_path = source_slide(tmp_path)
    project = Project(
        id="text-only-export",
        image_path=str(image_path),
        width=160,
        height=90,
        status="analyzed",
        components=[
            Component(
                id="title-text",
                type="text",
                bbox=BBox(x=20, y=17, width=50, height=7),
                text="Editable title",
                source="paddleocr",
            )
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    text_count, picture_count = shape_counts(pptx_path)
    assert text_count == 1
    assert picture_count == 0


def test_export_puts_detected_visuals_on_top_of_reconstruction_layer(tmp_path):
    image_path = source_slide(tmp_path)
    project = Project(
        id="visual-export",
        image_path=str(image_path),
        width=160,
        height=90,
        status="analyzed",
        components=[
            Component(
                id="red-icon",
                type="icon",
                bbox=BBox(x=105, y=45, width=30, height=30),
                source="sam3",
            )
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    text_count, picture_count = shape_counts(pptx_path)
    assert text_count == 0
    assert picture_count == 1

    with zipfile.ZipFile(pptx_path) as package:
        media_files = sorted(name for name in package.namelist() if name.startswith("ppt/media/"))
    assert len(media_files) == 1


def test_export_uses_residual_components_instead_of_full_slide_background(tmp_path):
    image_path = source_slide(tmp_path)
    project = Project(
        id="residual-export",
        image_path=str(image_path),
        width=160,
        height=90,
        status="analyzed",
        components=[
            Component(
                id="title-text",
                type="text",
                bbox=BBox(x=20, y=17, width=50, height=7),
                text="Editable title",
                source="paddleocr",
            ),
            Component(
                id="title-band",
                type="shape",
                bbox=BBox(x=12, y=12, width=136, height=18),
                source="opencv-residual",
            ),
            Component(
                id="red-icon",
                type="icon",
                bbox=BBox(x=105, y=45, width=30, height=30),
                source="sam3",
            ),
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    text_count, picture_count = shape_counts(pptx_path)
    assert text_count == 1
    assert picture_count == 1
    prs = Presentation(str(pptx_path))
    auto_shapes = [shape for shape in prs.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(auto_shapes) == 1


def test_export_preserves_visible_text_pixels_in_masked_visual_components(tmp_path):
    image_path = source_slide(tmp_path)
    mask_path = tmp_path / "title-band-mask.png"
    mask = Image.new("L", (160, 90), 0)
    draw = ImageDraw.Draw(mask)
    draw.rectangle([12, 12, 148, 30], fill=255)
    mask.save(mask_path)

    project = Project(
        id="masked-text-cleanup-export",
        image_path=str(image_path),
        width=160,
        height=90,
        status="analyzed",
        components=[
            Component(
                id="title-band",
                type="icon",
                bbox=BBox(x=12, y=12, width=136, height=18),
                mask_path=str(mask_path),
                source="sam3",
            ),
            Component(
                id="title-text",
                type="text",
                bbox=BBox(x=20, y=17, width=50, height=7),
                text="Editable title",
                source="paddleocr",
            ),
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    pictures = media_images(pptx_path)
    assert len(pictures) == 1
    red, green, blue, alpha = pictures[0].getpixel((20, 8))
    assert alpha > 0
    assert (red, green, blue) == (0, 0, 0)


def test_export_preserves_chart_pixels_under_editable_text(tmp_path):
    image_path = tmp_path / "chart-under-text.png"
    image = Image.new("RGB", (200, 120), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([30, 20, 170, 90], outline="#2563eb", width=2)
    draw.rectangle([50, 36, 120, 48], fill="black")
    image.save(image_path)
    project = Project(
        id="chart-under-text-export",
        image_path=str(image_path),
        width=200,
        height=120,
        status="analyzed",
        components=[
            Component(id="chart", type="chart", bbox=BBox(x=30, y=20, width=140, height=70), source="sam3"),
            Component(id="chart-label", type="text", bbox=BBox(x=50, y=36, width=70, height=12), text="Axis", source="paddleocr"),
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    pictures = media_images(pptx_path)
    assert len(pictures) == 1
    assert pictures[0].getpixel((25, 20))[:3] == (0, 0, 0)
    assert_text_shapes_are_frontmost(pptx_path)


def test_export_places_text_above_complex_backplate_picture(tmp_path):
    image_path = tmp_path / "complex-backplate.png"
    image = Image.new("RGB", (420, 140), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([24, 58, 390, 118], fill="#fdfdfd", outline="#bd110d", width=2)
    draw.polygon([(24, 58), (160, 58), (140, 118), (24, 118)], fill="#d71920")
    draw.rectangle([52, 72, 138, 100], fill="white")
    draw.rectangle([190, 78, 330, 94], fill="black")
    image.save(image_path)
    project = Project(
        id="complex-backplate-export",
        image_path=str(image_path),
        width=420,
        height=140,
        status="analyzed",
        components=[
            Component(id="banner", type="shape", bbox=BBox(x=24, y=58, width=366, height=60), source="opencv-residual"),
            Component(id="label", type="text", bbox=BBox(x=52, y=72, width=86, height=28), text="KEY TAKEAWAY", source="paddleocr"),
            Component(id="body", type="text", bbox=BBox(x=190, y=78, width=140, height=16), text="Editable summary", source="paddleocr"),
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    assert_text_shapes_are_frontmost(pptx_path)


def test_export_places_text_above_synthesized_info_card_shapes(tmp_path):
    image_path = tmp_path / "info-cards.png"
    image = Image.new("RGB", (360, 180), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([40, 92, 88, 138], fill="#fff8f5", outline="#d13a1a", width=2)
    draw.text((104, 96), "Strength", fill="#111111")
    draw.text((104, 118), "Event count", fill="#111111")
    image.save(image_path)
    project = Project(
        id="info-card-export",
        image_path=str(image_path),
        width=360,
        height=180,
        status="analyzed",
        components=[
            Component(id="icon-card", type="shape", bbox=BBox(x=40, y=92, width=48, height=46), source="opencv-residual"),
            Component(id="title", type="text", bbox=BBox(x=104, y=96, width=80, height=18), text="Strength", source="paddleocr"),
            Component(id="body", type="text", bbox=BBox(x=104, y=118, width=92, height=18), text="Event count", source="paddleocr"),
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    assert_text_shapes_are_frontmost(pptx_path)


def test_export_uses_source_crop_for_chart_assets_to_keep_annotations(tmp_path):
    image_path = tmp_path / "chart-source.png"
    image = Image.new("RGB", (140, 90), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([42, 8, 98, 22], fill="#d71920")
    draw.line([24, 66, 120, 36], fill="#2563eb", width=2)
    image.save(image_path)
    asset_path = tmp_path / "chart-asset.png"
    Image.new("RGBA", (140, 90), (255, 255, 255, 255)).save(asset_path)
    mask_path = tmp_path / "chart-mask.png"
    mask = Image.new("L", (140, 90), 0)
    mask_draw = ImageDraw.Draw(mask)
    mask_draw.line([24, 66, 120, 36], fill=255, width=8)
    mask.save(mask_path)
    project = Project(
        id="chart-source-crop-export",
        image_path=str(image_path),
        width=140,
        height=90,
        status="analyzed",
        components=[
            Component(
                id="chart",
                type="chart",
                bbox=BBox(x=0, y=0, width=140, height=90),
                asset_path=str(asset_path),
                mask_path=str(mask_path),
                source="sam3",
            )
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    pictures = media_images(pptx_path)
    assert len(pictures) == 1
    red, green, blue, alpha = pictures[0].getpixel((50, 14))
    assert alpha > 240
    assert (red, green, blue) == (215, 25, 32)


def test_export_writes_report_box_as_native_shape_not_picture(tmp_path):
    image_path = source_slide(tmp_path)
    project = Project(
        id="native-shape-export",
        image_path=str(image_path),
        width=160,
        height=90,
        status="analyzed",
        components=[
            Component(
                id="title-band",
                type="shape",
                bbox=BBox(x=12, y=12, width=136, height=18),
                source="opencv-residual",
            ),
            Component(
                id="title-text",
                type="text",
                bbox=BBox(x=20, y=17, width=50, height=7),
                text="Editable title",
                source="paddleocr",
            ),
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    auto_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(auto_shapes) == 1
    assert len(picture_shapes) == 0

    with zipfile.ZipFile(pptx_path) as package:
        xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
    assert '<a:srgbClr val="DBEAFE"/>' in xml


def test_export_places_text_above_synthesized_background_shapes(tmp_path):
    image_path = tmp_path / "dark-label.png"
    image = Image.new("RGB", (180, 100), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([20, 44, 150, 72], fill="#111827")
    draw.rectangle([38, 52, 120, 62], fill="white")
    image.save(image_path)
    project = Project(
        id="text-top-export",
        image_path=str(image_path),
        width=180,
        height=100,
        status="analyzed",
        components=[
            Component(
                id="label",
                type="text",
                bbox=BBox(x=38, y=52, width=82, height=10),
                text="Editable label",
                source="paddleocr",
            )
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    assert_text_shapes_are_frontmost(pptx_path)


def test_export_writes_arrow_component_as_connector_not_picture(tmp_path):
    image_path = source_slide(tmp_path)
    project = Project(
        id="native-arrow-export",
        image_path=str(image_path),
        width=160,
        height=90,
        status="analyzed",
        components=[
            Component(
                id="arrow",
                type="arrow",
                bbox=BBox(x=22, y=66, width=118, height=4),
                source="sam3",
            ),
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    line_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE]
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(line_shapes) == 1
    assert len(picture_shapes) == 0

    with zipfile.ZipFile(pptx_path) as package:
        xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<p:cxnSp>" in xml
    assert '<a:tailEnd type="triangle"/>' in xml


def test_export_uses_condensed_large_font_for_display_title(tmp_path):
    image_path = tmp_path / "display-title.png"
    Image.new("RGB", (1200, 400), "white").save(image_path)
    project = Project(
        id="display-title-export",
        image_path=str(image_path),
        width=1200,
        height=400,
        status="analyzed",
        components=[
            Component(
                id="title",
                type="text",
                bbox=BBox(x=20, y=8, width=840, height=92),
                text="Why Is One Punch Man So Strong?",
                source="paddleocr",
            )
        ],
    )

    pptx_path = export_pptx(project, project_store(tmp_path))

    prs = Presentation(str(pptx_path))
    paragraph_font = prs.slides[0].shapes[0].text_frame.paragraphs[0].font
    assert paragraph_font.name == "Impact"
    assert paragraph_font.bold
    assert paragraph_font.size.pt >= 34
