import io
import zipfile

from fastapi.testclient import TestClient
from PIL import Image, ImageDraw
from pptx import Presentation

from app.main import create_app
from app.settings import Settings


def make_slide_image() -> bytes:
    image = Image.new("RGB", (800, 450), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle([40, 35, 760, 95], outline="#173f73", width=3)
    draw.text((60, 55), "Editable title", fill="#111827")
    draw.rectangle([80, 135, 350, 335], outline="#1d4ed8", width=4)
    draw.line([105, 310, 160, 245, 230, 280, 315, 175], fill="#1d4ed8", width=5)
    draw.ellipse([560, 160, 635, 235], fill="#d71920")
    stream = io.BytesIO()
    image.save(stream, format="PNG")
    return stream.getvalue()


def client(tmp_path):
    settings = Settings(data_dir=tmp_path / "data")
    return TestClient(create_app(settings))


def upload_project(client):
    response = client.post(
        "/api/projects",
        files={"file": ("slide.png", make_slide_image(), "image/png")},
    )
    assert response.status_code == 200
    return response.json()


def test_upload_analyze_component_operations_and_export(tmp_path):
    test_client = client(tmp_path)
    project = upload_project(test_client)

    assert project["width"] == 800
    assert project["height"] == 450
    assert project["status"] == "uploaded"

    analyze = test_client.post(f"/api/projects/{project['id']}/analyze")
    assert analyze.status_code == 200
    analyzed = analyze.json()
    assert analyzed["status"] == "analyzed"
    assert len(analyzed["components"]) >= 2

    visible = [component for component in analyzed["components"] if not component["hidden"]]
    first_two = [visible[0]["id"], visible[1]["id"]]
    merge = test_client.patch(
        f"/api/projects/{project['id']}/components",
        json={"operation": "merge", "component_ids": first_two},
    )
    assert merge.status_code == 200
    merged_project = merge.json()
    assert any(component["source"] == "user-merge" for component in merged_project["components"])
    assert all(
        component["hidden"]
        for component in merged_project["components"]
        if component["id"] in first_two
    )

    delete_id = next(
        component["id"]
        for component in merged_project["components"]
        if not component["hidden"]
    )
    delete = test_client.patch(
        f"/api/projects/{project['id']}/components",
        json={"operation": "delete", "component_ids": [delete_id]},
    )
    assert delete.status_code == 200
    deleted_project = delete.json()
    assert next(component for component in deleted_project["components"] if component["id"] == delete_id)["hidden"]

    split = test_client.patch(
        f"/api/projects/{project['id']}/components",
        json={
            "operation": "split",
            "component_ids": [first_two[0]],
            "boxes": [
                {"x": 10, "y": 10, "width": 40, "height": 40},
                {"x": 60, "y": 10, "width": 40, "height": 40},
            ],
        },
    )
    assert split.status_code == 200
    split_project = split.json()
    assert sum(1 for component in split_project["components"] if component["source"] == "user-split") == 2

    export = test_client.post(f"/api/projects/{project['id']}/export/pptx")
    assert export.status_code == 200
    assert export.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    deck = tmp_path / "export.pptx"
    deck.write_bytes(export.content)
    prs = Presentation(str(deck))
    assert len(prs.slides) == 1
    text_shapes = [shape for shape in prs.slides[0].shapes if hasattr(shape, "text") and shape.text]
    picture_shapes = [shape for shape in prs.slides[0].shapes if shape.shape_type == 13]
    assert text_shapes or picture_shapes

    with zipfile.ZipFile(deck) as package:
        xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
    assert not (
        '<a:off x="0" y="0"/>' in xml
        and 'cy="6858000"' in xml
        and 'cx="1219' in xml
    )

    scene = test_client.get(f"/api/projects/{project['id']}/scene.svg")
    assert scene.status_code == 200
    assert scene.headers["content-type"].startswith("image/svg+xml")
    assert b"<svg" in scene.content
    assert b"data:image/png;base64," in scene.content or b"<text" in scene.content


def test_runtime_status_reports_missing_optional_ai_dependencies(tmp_path):
    test_client = client(tmp_path)
    response = test_client.get("/api/runtime")
    assert response.status_code == 200
    runtime = response.json()
    assert "paddleocr" in runtime
    assert "sam3" in runtime
    assert "torch" in runtime
