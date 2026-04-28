from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .image_editing import erase_regions
from .models import Project, SceneNode
from .scene import build_scene_graph, scene_paint_order
from .storage import ProjectStore

SLIDE_WIDTH = 13.333333
SLIDE_HEIGHT = 7.5
FONT = "Malgun Gothic"
SOURCE_CROP_VISUAL_TYPES = {"chart", "table", "diagram"}


def export_pptx(project: Project, store: ProjectStore) -> Path:
    output_path = store.export_dir(project.id) / "deck.pptx"
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    scale, offset_x, offset_y = _fit_to_slide(project.width, project.height)
    with Image.open(project.image_path) as opened:
        source = opened.convert("RGBA")
    scene = build_scene_graph(project, source)

    for node in scene_paint_order(scene.nodes):
        x = offset_x + node.bbox.x * scale
        y = offset_y + node.bbox.y * scale
        width = node.bbox.width * scale
        height = node.bbox.height * scale
        if node.kind == "text" and node.text:
            _add_text(slide, x, y, width, height, node.text, node.text_color)
        elif node.kind == "rect":
            _add_shape(slide, node, x, y, width, height)
        elif node.kind in {"line", "arrow"}:
            _add_connector(slide, node, scale, offset_x, offset_y)
        elif node.kind == "image":
            _add_picture(slide, source, node, x, y, width, height)

    _bring_text_shapes_to_front(slide)
    prs.save(output_path)
    return output_path


def _fit_to_slide(width: int, height: int) -> tuple[float, float, float]:
    scale = min(SLIDE_WIDTH / width, SLIDE_HEIGHT / height)
    content_w = width * scale
    content_h = height * scale
    return scale, (SLIDE_WIDTH - content_w) / 2, (SLIDE_HEIGHT - content_h) / 2


def _add_text(slide, x: float, y: float, width: float, height: float, text: str, text_color: str | None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(max(height, 0.22)))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = _font_name(width, height, text)
    paragraph.font.size = Pt(_font_size(width, height, text))
    paragraph.font.bold = _font_bold(width, height, text)
    paragraph.font.italic = _font_italic(text)
    paragraph.font.color.rgb = _rgb_color(text_color or "111827")
    return box


def _add_shape(slide, primitive: SceneNode, x: float, y: float, width: float, height: float):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    if primitive.fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb_color(primitive.fill_color)
    else:
        shape.fill.background()

    if primitive.line_color:
        shape.line.color.rgb = _rgb_color(primitive.line_color)
        if primitive.line_width:
            shape.line.width = Pt(primitive.line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_connector(slide, primitive: SceneNode, scale: float, offset_x: float, offset_y: float):
    x1 = offset_x + (primitive.x1 if primitive.x1 is not None else primitive.bbox.x) * scale
    y1 = offset_y + (primitive.y1 if primitive.y1 is not None else primitive.bbox.y) * scale
    x2 = offset_x + (primitive.x2 if primitive.x2 is not None else primitive.bbox.x + primitive.bbox.width) * scale
    y2 = offset_y + (primitive.y2 if primitive.y2 is not None else primitive.bbox.y + primitive.bbox.height) * scale
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    if primitive.line_color:
        connector.line.color.rgb = _rgb_color(primitive.line_color)
    if primitive.line_width:
        connector.line.width = Pt(primitive.line_width)
    if primitive.kind == "arrow":
        _add_tail_arrow(connector)
    return connector


def _add_tail_arrow(connector) -> None:
    line = connector._element.spPr.get_or_add_ln()
    for child in list(line):
        if child.tag.endswith("}tailEnd"):
            line.remove(child)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    line.append(tail)


def _add_picture(
    slide,
    source: Image.Image,
    primitive: SceneNode,
    x: float,
    y: float,
    width: float,
    height: float,
):
    asset = Path(primitive.asset_path) if primitive.asset_path else None
    if asset and asset.exists() and primitive.source_component_type not in SOURCE_CROP_VISUAL_TYPES:
        slide.shapes.add_picture(str(asset), Inches(x), Inches(y), width=Inches(width), height=Inches(height))
        return

    crop = source.crop(
        (
            max(0, round(primitive.bbox.x)),
            max(0, round(primitive.bbox.y)),
            min(source.width, round(primitive.bbox.x + primitive.bbox.width)),
            min(source.height, round(primitive.bbox.y + primitive.bbox.height)),
        )
    )
    if primitive.mask_path and Path(primitive.mask_path).exists() and primitive.source_component_type not in SOURCE_CROP_VISUAL_TYPES:
        crop = _apply_mask(crop, Path(primitive.mask_path), primitive)
    crop = erase_regions(crop, primitive.bbox, primitive.erase_boxes)

    _add_pil_picture(slide, crop, x, y, width, height)


def _add_pil_picture(slide, image: Image.Image, x: float, y: float, width: float, height: float):
    stream = BytesIO()
    image.save(stream, format="PNG")
    stream.seek(0)
    slide.shapes.add_picture(stream, Inches(x), Inches(y), width=Inches(width), height=Inches(height))


def _bring_text_shapes_to_front(slide) -> None:
    text_elements = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if not shape.text.strip():
            continue
        text_elements.append(shape._element)
    for element in text_elements:
        slide.shapes._spTree.remove(element)
        slide.shapes._spTree.append(element)


def _apply_mask(crop: Image.Image, mask_path: Path, primitive: SceneNode) -> Image.Image:
    mask = Image.open(mask_path).convert("L")
    mask_crop = mask.crop(
        (
            max(0, round(primitive.bbox.x)),
            max(0, round(primitive.bbox.y)),
            min(mask.width, round(primitive.bbox.x + primitive.bbox.width)),
            min(mask.height, round(primitive.bbox.y + primitive.bbox.height)),
        )
    ).resize(crop.size)
    crop = crop.convert("RGBA")
    crop.putalpha(mask_crop)
    return crop


def _font_size(width: float, height: float, text: str) -> float:
    longest = max(text.splitlines() or [text], key=len)
    if _is_section_label_text(text):
        by_width = width * 72 / max(_text_units(longest, latin_width=0.5), 1.0) * 0.96
        by_height = height * 72 * 0.78
        return max(10.0, min(48.0, by_width, by_height))
    if _is_latin_display_text(width, height, text):
        by_width = width * 72 / max(_text_units(longest, latin_width=0.36), 1.0) * 0.96
        by_height = height * 72 / max(text.count("\n") + 1, 1) * 0.86
        return max(10.0, min(72.0, by_width, by_height))
    by_width = width * 72 / max(_text_units(longest), 1.0) * 0.98
    by_height = height * 72 / max(text.count("\n") + 1, 1) * 0.75
    return max(7.0, min(34.0, by_width, by_height))


def _font_name(width: float, height: float, text: str) -> str:
    if _is_latin_display_text(width, height, text):
        return "Impact"
    if _is_latin_text(text):
        return "Arial"
    return FONT


def _font_bold(width: float, height: float, text: str) -> bool:
    return _is_latin_display_text(width, height, text) or _is_short_uppercase_label(text)


def _font_italic(text: str) -> bool:
    stripped = text.strip()
    return stripped.startswith(("Narrative,", "A narrative", "The narrative"))


def _text_units(text: str, latin_width: float = 0.6) -> float:
    total = 0.0
    for char in text:
        code = ord(char)
        if 0xAC00 <= code <= 0xD7A3 or 0x4E00 <= code <= 0x9FFF:
            total += 1.0
        elif char.isspace():
            total += 0.35
        else:
            total += latin_width
    return total


def _is_latin_display_text(width: float, height: float, text: str) -> bool:
    stripped = text.strip()
    return (
        height >= 0.36
        and len(stripped) >= 12
        and _latin_character_ratio(stripped) >= 0.72
        and any(char.isupper() for char in stripped)
    )


def _is_latin_text(text: str) -> bool:
    return _latin_character_ratio(text.strip()) >= 0.65


def _is_short_uppercase_label(text: str) -> bool:
    letters = [char for char in text if char.isalpha()]
    if not letters:
        return False
    uppercase = sum(1 for char in letters if char.upper() == char)
    return len(text.strip()) <= 56 and uppercase / max(1, len(letters)) >= 0.75


def _is_section_label_text(text: str) -> bool:
    stripped = text.strip()
    return (
        _is_short_uppercase_label(stripped)
        and len(stripped) >= 2
        and (stripped[1] in {")", "]"} or stripped[:3].upper() == "VS.")
    )


def _latin_character_ratio(text: str) -> float:
    non_space = [char for char in text if not char.isspace()]
    if not non_space:
        return 0.0
    latin = sum(1 for char in non_space if ord(char) < 128)
    return latin / len(non_space)


def _rgb_color(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
